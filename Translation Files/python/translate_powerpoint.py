#-----------------------------------------------------
#概要：PowerPointを翻訳するプログラム
#補足：複数スライドを翻訳する
#　　　本Pythonファイルと同じフォルダに翻訳元PowerPointを格納しておく
#　　　翻訳後のファイル名は「翻訳元ファイル名(Translated).拡張子」とする
#-----------------------------------------------------

from openpyxl.descriptors.base import String
from pptx import Presentation
from time import sleep
import os
from googletrans import Translator

#-----------------------------------------------------
#翻訳実行メソッド
#-----------------------------------------------------
def translate(trans_obj, file_path: String, src_lang_code: String, dest_lang_code: String):
    trans_obj.updateLog('Start translation.\n')
    
    prs = Presentation(file_path)
    translator = Translator() #翻訳準備
    total_sheet_num = len(prs.slides) #スライド数取得

    for ns, slide in enumerate(prs.slides):
        trans_obj.updateLog('Translating ' + str(ns + 1) + "/" + str(total_sheet_num) + " sheet...\n") #翻訳しているページ数を出力

        for nsh, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            for np, paragraph in enumerate(shape.text_frame.paragraphs):
                for rs, run in enumerate(paragraph.runs):
                    str_in = run.text
                    str_out = translator.translate(str_in, src=src_lang_code, dest=dest_lang_code).text
                    prs.slides[ns].shapes[nsh].text_frame.paragraphs[np].runs[rs].text = str_out
                    sleep(1.5) #サーバ負荷軽減処理。

    file_name_split = os.path.splitext(os.path.basename(file_path)) #翻訳元ファイル名を拡張子以外と拡張子に分割
    prs.save(file_name_split[0] + "(Translated)" + file_name_split[1])
    trans_obj.updateLog('Translation was completed.\n') #翻訳終了メッセージ
